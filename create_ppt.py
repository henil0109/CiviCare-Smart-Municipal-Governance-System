from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "CiviCare"
    subtitle.text = "Smart Municipal Governance System\n\nDeveloped by:\nHENIL PATEL (System Architect, AI & Full Stack)\nMEET PATEL (Frontend, UI/UX, Documentation)"

    # Helper function to add a standard slide
    def add_slide(title_text, content_text):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        tf.text = content_text
        return slide

    # 2. Objective & Need
    add_slide(
        "Objective & Need",
        "• The Need: Rapid urban growth leads to civic issues (broken infrastructure, waste mismanagement).\n"
        "• Traditional systems are slow, paper-based, and lack transparency.\n"
        "• The Solution (CiviCare): Bridge the gap between citizens and municipal authorities.\n"
        "• Uses AI to automate triaging of civic complaints.\n"
        "• Provides gamified, transparent, real-time tracking for citizens."
    )

    # 3. Scope & Limitations
    add_slide(
        "Scope & Limitations",
        "Scope:\n"
        "• Public portal with identity verification & evidence-based reporting.\n"
        "• AI engine for automatic categorization, priority, and risk assessment.\n"
        "• Hierarchical dashboard for Admins, Supervisors, and Field Teams.\n"
        "• Gamification system (XP/Badges) to reward active citizens.\n\n"
        "Limitations:\n"
        "• AI relies on heuristic weighting rather than continuous neural learning.\n"
        "• Live GPS tracking is currently simulated.\n"
        "• Requires internet access, limiting non-digital users."
    )

    # 4. Technology Stack
    add_slide(
        "Technology Stack",
        "Frontend:\n"
        "• React.js (Vite), Tailwind CSS, Framer Motion\n\n"
        "Backend:\n"
        "• Python (Flask), JWT Authentication, Werkzeug\n"
        "• Brevo HTTP API for transactional emails\n\n"
        "Database:\n"
        "• MongoDB (Atlas) via PyMongo"
    )

    # 5. Core Features & Concepts
    add_slide(
        "Core Features & Concepts",
        "• Automated AI Triage: Instantly routes issues to the right department.\n"
        "• Role-Based Access (RBAC): Dedicated views for Admin, Supervisor, Citizen.\n"
        "• CiviBot AI Assistant: 24/7 floating widget for FAQs and guidance.\n"
        "• Civic Gamification: XP & Levels for constructive civic duty.\n"
        "• Asynchronous Architecture: Email dispatching without UI blocking.\n"
        "• Modern UX/UI: Glassmorphism, smooth animations, dark/light modes."
    )

    # 6. System Workflow
    add_slide(
        "System Workflow",
        "1. User Onboarding: Async email verification and JWT token access.\n"
        "2. Issue Reporting: AI analyzes text for category, priority, and safety risk.\n"
        "3. Delegation: Admins assign AI-triaged tasks to Supervisors.\n"
        "4. Field Work: Supervisors dispatch teams & update live timeline.\n"
        "5. Resolution & Rewards: Final AI report generated, citizen earns XP/Badges."
    )

    # 7. System Design & Architecture
    add_slide(
        "System Design & Architecture",
        "• Client-Server Model: React frontend independent from Flask backend.\n"
        "• RESTful JSON APIs for seamless communication.\n"
        "• Stateless Backend (JWT): Allows horizontal scaling.\n"
        "• Security: Password hashing, CORS management, secure token exchange."
    )

    # 8. Results & Impact
    add_slide(
        "Results & Impact",
        "• Faster Response Times: Emergencies are instantly flagged by AI.\n"
        "• Higher Accountability: Every step is timestamped in a public timeline.\n"
        "• Resource Optimization: AI estimates required team size and costs.\n"
        "• Engaged Citizens: Transforms complaining into a rewarding civic duty."
    )

    # Save presentation
    prs.save("CiviCare_Presentation.pptx")
    print("Presentation successfully created at CiviCare_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
